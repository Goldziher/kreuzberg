"""This module provides functions to extract textual content from files.

It includes vendored code:

- The extract PPTX logic is based on code vendored from `markitdown` to extract text from PPTX files.
    See: https://github.com/microsoft/markitdown/blob/main/src/markitdown/_markitdown.py
    Refer to the markitdown repository for it's license (MIT).
"""

from __future__ import annotations

import gc
import re

try:
    import defusedxml.ElementTree as ET  # noqa: N817
except ImportError:
    import xml.etree.ElementTree as ET
from contextlib import suppress
from html import escape
from io import BytesIO
from pathlib import Path
from typing import TYPE_CHECKING, ClassVar
from zipfile import ZipFile

import pptx
from anyio import Path as AsyncPath
from pptx.enum.shapes import MSO_SHAPE_TYPE

from kreuzberg._extractors._base import Extractor
from kreuzberg._mime_types import MARKDOWN_MIME_TYPE, POWER_POINT_MIME_TYPE
from kreuzberg._types import ExtractionResult
from kreuzberg._utils._string import normalize_spaces

if TYPE_CHECKING:  # pragma: no cover
    from pptx.presentation import Presentation

    from kreuzberg._types import Metadata


class PresentationExtractor(Extractor):
    """Extractor for PowerPoint (.pptx) files.

    This extractor processes PowerPoint presentations and converts their content into Markdown format.
    It handles slides, shapes, images, tables, and slide notes, preserving the structure and content
    of the presentation in a readable text format.

    The extractor provides both synchronous and asynchronous methods for processing files either
    from disk or from bytes in memory.
    """

    SUPPORTED_MIME_TYPES: ClassVar[set[str]] = {POWER_POINT_MIME_TYPE}

    async def extract_bytes_async(self, content: bytes) -> ExtractionResult:
        """Asynchronously extract content from PowerPoint file bytes.

        Args:
            content: Raw bytes of the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.
        """
        return self._extract_pptx(content)

    async def extract_path_async(self, path: Path) -> ExtractionResult:
        """Asynchronously extract content from a PowerPoint file on disk.

        Args:
            path: Path to the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.
        """
        content = await AsyncPath(path).read_bytes()
        return self._extract_pptx(content)

    def extract_bytes_sync(self, content: bytes) -> ExtractionResult:
        """Synchronously extract content from PowerPoint file bytes.

        Args:
            content: Raw bytes of the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.
        """
        return self._extract_pptx(content)

    def extract_path_sync(self, path: Path) -> ExtractionResult:
        """Synchronously extract content from a PowerPoint file on disk.

        Args:
            path: Path to the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.
        """
        content = Path(path).read_bytes()
        return self._extract_pptx(content)

    def _extract_pptx(self, file_contents: bytes) -> ExtractionResult:
        """Process PowerPoint file contents and convert to Markdown.

        This method handles the core logic of extracting content from a PowerPoint file.
        It processes:
        - Slide titles and content
        - Images (with alt text if available)
        - Tables (converted to HTML format)
        - Text frames
        - Slide notes

        Args:
            file_contents: Raw bytes of the PowerPoint file to process.

        Returns:
            ExtractionResult: Contains the extracted content in Markdown format,
                the MIME type, and any additional metadata.

        Notes:
            The extraction preserves the following elements:
            - Slide numbers (as HTML comments)
            - Images (converted to Markdown image syntax with alt text)
            - Tables (converted to HTML table syntax)
            - Text content (with titles properly formatted)
            - Slide notes (under a dedicated section for each slide)
        """
        md_content = []  # Use list for better memory efficiency
        presentation = pptx.Presentation(BytesIO(file_contents))

        # Process metadata first before slides to reduce peak memory
        metadata = self._extract_presentation_metadata(presentation)

        # Get total slide count for memory management
        total_slides = len(presentation.slides)

        for index, slide in enumerate(presentation.slides):
            slide_content = [f"\n\n<!-- Slide number: {index + 1} -->\n"]

            title = None
            if hasattr(slide.shapes, "title"):
                title = slide.shapes.title

            for shape in slide.shapes:
                if not hasattr(shape, "shape_type"):
                    continue

                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE or (
                    shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and hasattr(shape, "image")
                ):
                    alt_text = ""
                    with suppress(AttributeError):
                        alt_text = shape._element._nvXxPr.cNvPr.attrib.get("descr", "")  # noqa: SLF001

                    filename = re.sub(r"\W", "", shape.name) + ".jpg"
                    slide_content.append(f"\n![{alt_text if alt_text else shape.name}]({filename})\n")

                elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    html_table = "<table>"
                    first_row = True

                    for row in shape.table.rows:
                        html_table += "<tr>"

                        for cell in row.cells:
                            tag = "th" if first_row else "td"
                            html_table += f"<{tag}>{escape(cell.text)}</{tag}>"

                        html_table += "</tr>"
                        first_row = False

                    html_table += "</table>"
                    slide_content.append("\n" + html_table + "\n")

                elif shape.has_text_frame:
                    slide_content.append("# " + shape.text.lstrip() + "\n" if shape == title else shape.text + "\n")

            # Join slide content and handle notes
            if slide.has_notes_slide:
                slide_content.append("\n\n### Notes:\n")
                notes_frame = slide.notes_slide.notes_text_frame

                if notes_frame is not None:  # pragma: no branch
                    slide_content.append(notes_frame.text)

            # Add slide content to main content
            md_content.append("".join(slide_content).strip())

            # Force garbage collection every 10 slides for large presentations
            large_presentation_threshold = 50
            if total_slides > large_presentation_threshold and (index + 1) % 10 == 0:
                gc.collect(0)

        # Clear presentation object to free memory
        del presentation
        gc.collect()

        return ExtractionResult(
            content=normalize_spaces("".join(md_content)),
            mime_type=MARKDOWN_MIME_TYPE,
            metadata=metadata,
            chunks=[],
        )

    @staticmethod
    def _extract_presentation_metadata(presentation: Presentation) -> Metadata:
        """Extract metadata from a presentation instance.

        Args:
            presentation: A `Presentation` object representing the PowerPoint file.

        Returns:
            PresentationMetadata: Object containing presentation-specific metadata fields.
        """
        metadata: Metadata = {}

        for metadata_key, core_property_key in [
            ("authors", "author"),
            ("comments", "comments"),
            ("status", "content_status"),
            ("created_by", "created"),
            ("identifier", "identifier"),
            ("keywords", "keywords"),
            ("modified_by", "last_modified_by"),
            ("modified_at", "modified"),
            ("version", "revision"),
            ("subject", "subject"),
            ("title", "title"),
            ("version", "version"),
        ]:
            if core_property := getattr(presentation.core_properties, core_property_key, None):
                metadata[metadata_key] = core_property  # type: ignore[literal-required]

        if presentation.core_properties.language:
            metadata["languages"] = [presentation.core_properties.language]

        if presentation.core_properties.category:
            metadata["categories"] = [presentation.core_properties.category]

        fonts = set()
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not hasattr(shape, "text_frame"):
                    continue

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if hasattr(run, "font") and run.font.name:
                            fonts.add(run.font.name)

        if fonts:
            metadata["fonts"] = list(fonts)

        return metadata

    def _extract_pptx_memory_efficient(self, content: bytes) -> ExtractionResult:
        """Memory-efficient PPTX extraction that avoids loading images.

        This method processes PPTX files as ZIP archives and only loads
        the XML content we need, avoiding expensive image loading.

        Args:
            content: Raw bytes of the PPTX file

        Returns:
            ExtractionResult with extracted text content
        """
        md_content = ""
        metadata: Metadata = {}

        try:
            with ZipFile(BytesIO(content), "r") as zip_file:
                # Get list of slide XML files
                slide_files = [
                    f for f in zip_file.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")
                ]
                slide_files.sort()  # Ensure proper order

                # Extract metadata from core properties
                if "docProps/core.xml" in zip_file.namelist():
                    with zip_file.open("docProps/core.xml") as core_file:
                        metadata.update(self._extract_core_properties(core_file.read()))

                # Process each slide
                for i, slide_file in enumerate(slide_files, 1):
                    md_content += f"\n\n<!-- Slide number: {i} -->\n"

                    with zip_file.open(slide_file) as slide_xml:
                        slide_content = self._extract_slide_text_from_xml(slide_xml.read())
                        md_content += slide_content

                    # Also check for notes
                    notes_file = slide_file.replace("/slide", "/notesSlide").replace(".xml", "Notes.xml")
                    if notes_file in zip_file.namelist():
                        with zip_file.open(notes_file) as notes_xml:
                            notes_content = self._extract_notes_from_xml(notes_xml.read())
                            if notes_content.strip():
                                md_content += f"\n\n### Notes:\n{notes_content}"

        except Exception:  # noqa: BLE001
            # Fallback to regular extraction if ZIP parsing fails
            return self._extract_pptx(content)

        return ExtractionResult(
            content=normalize_spaces(md_content),
            mime_type=MARKDOWN_MIME_TYPE,
            metadata=metadata,
            chunks=[],
        )

    def _extract_core_properties(self, xml_content: bytes) -> Metadata:
        """Extract metadata from core.xml."""
        metadata: Metadata = {}
        try:
            root = ET.fromstring(xml_content)  # noqa: S314

            # Define namespace mappings
            namespaces = {
                "cp": "http://schemas.openxmlformats.org/package/2006/metadata/core-properties",
                "dc": "http://purl.org/dc/elements/1.1/",
                "dcterms": "http://purl.org/dc/terms/",
            }

            # Extract common properties
            title_elem = root.find(".//dc:title", namespaces)
            if title_elem is not None and title_elem.text:
                metadata["title"] = title_elem.text

            creator_elem = root.find(".//dc:creator", namespaces)
            if creator_elem is not None and creator_elem.text:
                metadata["authors"] = [creator_elem.text]

            description_elem = root.find(".//dc:description", namespaces)
            if description_elem is not None and description_elem.text:
                metadata["description"] = description_elem.text

        except ET.ParseError:
            pass  # Ignore XML parsing errors

        return metadata

    def _extract_slide_text_from_xml(self, xml_content: bytes) -> str:
        """Extract text content from slide XML."""
        try:
            root = ET.fromstring(xml_content)  # noqa: S314

            # Find all text elements (t tags in the XML)
            text_elements = root.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t")

            slide_text = ""
            for elem in text_elements:
                if elem.text:
                    slide_text += elem.text + " "

            return slide_text.strip()

        except ET.ParseError:
            return ""

    def _extract_notes_from_xml(self, xml_content: bytes) -> str:
        """Extract notes content from notes slide XML."""
        try:
            root = ET.fromstring(xml_content)  # noqa: S314

            # Find all text elements in notes
            text_elements = root.findall(".//{http://schemas.openxmlformats.org/drawingml/2006/main}t")

            notes_text = ""
            for elem in text_elements:
                if elem.text:
                    notes_text += elem.text + " "

            return notes_text.strip()

        except ET.ParseError:
            return ""
